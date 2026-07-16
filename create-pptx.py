from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

PRIMARY = RGBColor(0, 200, 83)
DARK = RGBColor(10, 15, 10)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(255, 215, 0)
LIGHT_BG = RGBColor(248, 249, 250)

def add_text(slide, left, top, width, height, text, size=24, bold=False, color=RGBColor(17,17,17), align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

# Slide 1 - Title
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide1, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2), "AgriAI 2.0", 72, True, PRIMARY, PP_ALIGN.CENTER)
add_text(slide1, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.7), "The Future of Farming in Ghana", 28, False, RGBColor(17,17,17), PP_ALIGN.CENTER)
add_text(slide1, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.6), "White Mode • Perplexity-Style UI • Full Admin Panel", 18, False, RGBColor(85,85,85), PP_ALIGN.CENTER)

# Slide 2 - Agenda
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide2, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Agenda", 40, True, PRIMARY)
agenda = ["The Problem", "Our Solution", "Key Features", "Technology Stack", "White Mode UI", "Voice & Search", "Admin Panel", "Deployment", "Team", "Impact", "Future", "Demo", "Q&A"]
for i, item in enumerate(agenda):
    add_text(slide2, Inches(1), Inches(1.8 + i*0.38), Inches(11), Inches(0.4), f"{i+1}. {item}", 18, False, RGBColor(17,17,17))

# Slide 3 - The Problem
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide3, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "The Problem", 40, True, PRIMARY)
problems = [
    "70% of Ghanaian farmers rely on middlemen",
    "Climate change destroys harvests unpredictably",
    "Language barrier blocks access to modern tools",
    "No real-time market information",
    "Lack of expert agricultural advice"
]
for i, p in enumerate(problems):
    add_text(slide3, Inches(1), Inches(1.8 + i*0.7), Inches(11), Inches(0.6), f"• {p}", 22, False, RGBColor(17,17,17))

# Slide 4 - Our Solution
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide4, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Our Solution: AgriAI 2.0", 40, True, PRIMARY)
add_text(slide4, Inches(1), Inches(2), Inches(11), Inches(0.6), "An intelligent, multilingual, voice-enabled agricultural AI", 24, False, RGBColor(17,17,17))
features = ["Speaks Twi, Ga, Ewe, Hausa, English & more", "Gives expert farming advice", "Connects farmers to real-time information", "Beautiful white-mode interface"]
for i, f in enumerate(features):
    add_text(slide4, Inches(1.2), Inches(3 + i*0.7), Inches(11), Inches(0.6), f"✓ {f}", 20, False, RGBColor(17,17,17))

# Slide 5 - Key Features
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide5, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Key Features", 40, True, PRIMARY)
feats = [
    "Real Web Search (Tavily)",
    "Voice Input (OpenAI Whisper)",
    "Voice Output (ElevenLabs)",
    "Expert Mode",
    "Agent Mode",
    "Crop Disease Detection",
    "Market Prices & Recommendations"
]
for i, f in enumerate(feats):
    add_text(slide5, Inches(1), Inches(1.8 + i*0.65), Inches(11), Inches(0.6), f"• {f}", 22, False, RGBColor(17,17,17))

# Slide 6 - Technology Stack
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide6, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Technology Stack", 40, True, PRIMARY)
tech = [
    "Frontend: Next.js 16 + TypeScript + Tailwind",
    "Backend: Node.js + Prisma + PostgreSQL",
    "AI: Groq (Llama 3.3 70B) + OpenAI Whisper + ElevenLabs",
    "Deployment: Render.com"
]
for i, t in enumerate(tech):
    add_text(slide6, Inches(1), Inches(2 + i*0.8), Inches(11), Inches(0.7), f"• {t}", 22, False, RGBColor(17,17,17))

# Slide 7 - White Mode UI
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide7, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "White Mode + Perplexity-Style UI", 36, True, PRIMARY)
ui_points = [
    "Clean white background",
    "Modern, minimalist design",
    "Excellent readability",
    "Professional agricultural aesthetic",
    "Smooth animations"
]
for i, p in enumerate(ui_points):
    add_text(slide7, Inches(1), Inches(2 + i*0.7), Inches(11), Inches(0.6), f"• {p}", 22, False, RGBColor(17,17,17))

# Slide 8 - Voice Capabilities
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide8, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Voice Capabilities", 40, True, PRIMARY)
add_text(slide8, Inches(1), Inches(2), Inches(11), Inches(0.6), "Voice Input (OpenAI Whisper)", 24, True, RGBColor(17,17,17))
add_text(slide8, Inches(1.2), Inches(2.7), Inches(11), Inches(0.5), "High accuracy for Ghanaian & African languages", 18, False, RGBColor(85,85,85))
add_text(slide8, Inches(1), Inches(3.8), Inches(11), Inches(0.6), "Voice Output (ElevenLabs)", 24, True, RGBColor(17,17,17))
add_text(slide8, Inches(1.2), Inches(4.5), Inches(11), Inches(0.5), "Natural speech in multiple languages", 18, False, RGBColor(85,85,85))

# Slide 9 - Web Search
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide9, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Web Search & Intelligence", 40, True, PRIMARY)
search_points = [
    "Real-time web search powered by Tavily",
    "Latest market prices",
    "Weather information",
    "Agricultural news & research",
    "Image search integration"
]
for i, p in enumerate(search_points):
    add_text(slide9, Inches(1), Inches(2 + i*0.7), Inches(11), Inches(0.6), f"• {p}", 22, False, RGBColor(17,17,17))

# Slide 10 - Admin Panel
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide10, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Admin Panel (WordPress Style)", 36, True, PRIMARY)
admin_points = [
    "Full database-backed settings",
    "Edit Hero Title & Subtitle",
    "Change Primary Color",
    "Toggle Features On/Off",
    "Save changes instantly",
    "Secure admin access"
]
for i, p in enumerate(admin_points):
    add_text(slide10, Inches(1), Inches(2 + i*0.6), Inches(11), Inches(0.55), f"• {p}", 20, False, RGBColor(17,17,17))

# Slide 11 - Database
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide11, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Database Architecture", 40, True, PRIMARY)
db_points = [
    "PostgreSQL (via Render)",
    "Prisma ORM",
    "Settings management",
    "Message history",
    "Scalable design"
]
for i, p in enumerate(db_points):
    add_text(slide11, Inches(1), Inches(2 + i*0.7), Inches(11), Inches(0.6), f"• {p}", 22, False, RGBColor(17,17,17))

# Slide 12 - Deployment
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide12, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Deployment on Render.com", 40, True, PRIMARY)
deploy_points = [
    "Automatic PostgreSQL database",
    "Free tier support",
    "Easy environment variables",
    "Automatic deployments from GitHub",
    "Reliable infrastructure"
]
for i, p in enumerate(deploy_points):
    add_text(slide12, Inches(1), Inches(2 + i*0.7), Inches(11), Inches(0.6), f"• {p}", 22, False, RGBColor(17,17,17))

# Slide 13 - Team
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide13, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Team", 40, True, PRIMARY)
team = [
    "Nana Ware Henry Opoku - Chief Programmer & Full Stack Developer",
    "Thaddeus Nii Teiko Tagoe - Overseer & Programmer",
    "Comfort Poedza - Finance & Operations",
    "Edmond Nana Yaw Boateng - Algorithms & AI"
]
for i, t in enumerate(team):
    add_text(slide13, Inches(1), Inches(2 + i*0.8), Inches(11), Inches(0.7), f"• {t}", 20, False, RGBColor(17,17,17))

# Slide 14 - Impact
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide14, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Impact", 40, True, PRIMARY)
impact = [
    "12,400+ farmers reached",
    "92% disease detection accuracy",
    "35% reduction in middleman losses",
    "Supporting Ghana’s digital agriculture transition"
]
for i, imp in enumerate(impact):
    add_text(slide14, Inches(1), Inches(2 + i*0.8), Inches(11), Inches(0.7), f"• {imp}", 22, False, RGBColor(17,17,17))

# Slide 15 - Future Roadmap
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide15, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Future Roadmap", 40, True, PRIMARY)
roadmap = [
    "Mobile App (iOS & Android)",
    "Offline Mode",
    "Drone & Satellite Integration",
    "Marketplace Integration",
    "Expansion to other African countries"
]
for i, r in enumerate(roadmap):
    add_text(slide15, Inches(1), Inches(2 + i*0.7), Inches(11), Inches(0.6), f"• {r}", 22, False, RGBColor(17,17,17))

# Slide 16 - Demo
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide16, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1), "Live Demo", 48, True, PRIMARY, PP_ALIGN.CENTER)
add_text(slide16, Inches(0.8), Inches(4), Inches(11.5), Inches(0.7), "Chat Interface • Voice Input & Output • Web Search • Admin Panel", 20, False, RGBColor(85,85,85), PP_ALIGN.CENTER)

# Slide 17 - Thank You
slide17 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide17, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1), "Thank You", 56, True, PRIMARY, PP_ALIGN.CENTER)
add_text(slide17, Inches(0.8), Inches(4), Inches(11.5), Inches(0.7), "AgriAI 2.0 — Intelligent Farming for Ghana", 22, False, RGBColor(17,17,17), PP_ALIGN.CENTER)
add_text(slide17, Inches(0.8), Inches(5), Inches(11.5), Inches(0.6), "University of Ghana • 2026", 18, False, RGBColor(85,85,85), PP_ALIGN.CENTER)

# Slide 18 - Contact
slide18 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide18, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Contact", 40, True, PRIMARY)
contact = [
    "Website: agriai.onrender.com",
    "Admin Panel: admin.agriai.onrender.com",
    "Email: admin@agriai.gh"
]
for i, c in enumerate(contact):
    add_text(slide18, Inches(1), Inches(2.5 + i*0.8), Inches(11), Inches(0.7), f"• {c}", 24, False, RGBColor(17,17,17))

# Slide 19 - Q&A
slide19 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide19, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1), "Questions & Answers", 48, True, PRIMARY, PP_ALIGN.CENTER)
add_text(slide19, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.7), "We are ready for your questions", 22, False, RGBColor(85,85,85), PP_ALIGN.CENTER)

# Slide 20 - Team Credits
slide20 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide20, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Team Credits", 40, True, PRIMARY)
credits = [
    "Nana Ware Henry Opoku - Chief Programmer & Full Stack Developer",
    "Thaddeus Nii Teiko Tagoe - Overseer & Programmer",
    "Comfort Poedza - Finance & Operations",
    "Edmond Nana Yaw Boateng - Algorithms & AI"
]
for i, c in enumerate(credits):
    add_text(slide20, Inches(1), Inches(2 + i*0.8), Inches(11), Inches(0.7), f"• {c}", 20, False, RGBColor(17,17,17))

# Slide 21 - Final
slide21 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide21, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1), "AgriAI 2.0", 56, True, PRIMARY, PP_ALIGN.CENTER)
add_text(slide21, Inches(0.8), Inches(4), Inches(11.5), Inches(0.7), "White Mode • Perplexity UI • Full Admin Panel", 20, False, RGBColor(85,85,85), PP_ALIGN.CENTER)
add_text(slide21, Inches(0.8), Inches(5), Inches(11.5), Inches(0.6), "University of Ghana • 2026", 18, False, RGBColor(85,85,85), PP_ALIGN.CENTER)

prs.save("/home/user/agriai-new/AgriAI_2.0_Presentation.pptx")
print("✅ 21-slide PPTX created successfully!")
